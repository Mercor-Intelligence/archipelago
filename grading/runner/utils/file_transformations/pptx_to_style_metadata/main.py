"""Extract font, color, layout, and shape metadata from a PPTX file.

The dict-returning extractor is a domain-specific utility used by the
pptx_style_verifier eval. pptx_to_style_metadata_output wraps it as a
registered transformation for the generic multi-representation judge, emitting
a bounded per-slide style summary rather than the full dict.

See ../README.md for the rules these extractors share: resolve references to
the value that renders, name the source, never invent one, and cap every
section with the cap disclosed.
"""

import asyncio
from collections import Counter
from io import BytesIO
from typing import Any

from loguru import logger
from pptx import Presentation
from pptx.util import Emu

from ..models import TransformationOutput
from ..style_metadata_cache import cached_style_text, is_ooxml_package
from ..xml_utils import DML_NS, PML_NS, xml_escape


def _emu_to_pt(emu_val: int | Emu | None) -> float | None:
    """Convert EMU to points. Returns None if input is None."""
    if emu_val is None:
        return None
    return round(int(emu_val) / 12700, 1)


def _rgb_to_hex(rgb: Any) -> str | None:
    """Convert an RGBColor to a hex string like '#RRGGBB'. Returns None if not set."""
    if rgb is None:
        return None
    try:
        return f"#{rgb}"
    except Exception:
        return None


def _theme_color_name(color: Any) -> str | None:
    """A theme colour as "accent1" / "dark1", plus any brightness tweak.

    python-pptx raises "no .rgb property on color type '_SchemeColor'" for a run
    coloured from the theme, so catching AttributeError and moving on left no
    colour information at all — and colouring by theme is what PowerPoint
    templates do by default. The theme slot is the honest answer: the concrete
    RGB lives in the theme part and can be overridden per layout, so naming the
    slot is both correct and stable.

    Brightness carries PowerPoint's "Lighter 40% / Darker 25%" variants, without
    which two visibly different shades collapse to the same slot.
    """
    try:
        theme = color.theme_color
    except Exception:
        return None
    # Falsy, not None: python-pptx returns the NOT_THEME_COLOR sentinel (an
    # int enum of value 0) for every non-scheme colour, and every real member
    # is nonzero — so `not theme` screens the sentinel and None alike.
    if not theme:
        return None
    # MSO_THEME_COLOR_INDEX.ACCENT_1 -> "accent1"
    name = str(theme).split()[0].lower().replace("_", "")
    try:
        brightness = color.brightness
    except Exception:
        brightness = 0
    if brightness:
        return f"{name}{brightness:+.0%}".replace("%", "pct")
    return name


def _extract_run_style(run: Any) -> dict[str, Any]:
    """Extract style properties from a single text run."""
    font = run.font
    color = font.color
    color_rgb = None
    color_theme = None
    if color is not None:
        try:
            color_rgb = _rgb_to_hex(color.rgb)
        except AttributeError:
            # Not an explicit RGB colour. A theme reference is still real colour
            # information, so resolve the slot rather than dropping it.
            color_theme = _theme_color_name(color)
    return {
        "text": run.text,
        "font_name": font.name,  # None means inherited from slide master
        "font_size_pt": _emu_to_pt(font.size),
        "font_color_rgb": color_rgb,
        # Additive: font_color_rgb keeps its meaning (explicit hex, or None), so
        # existing consumers see a new key rather than a changed value.
        "font_color_theme": color_theme,
        "bold": font.bold,  # None means inherited
        "italic": font.italic,  # None means inherited
    }


def _shape_type_name(shape: Any) -> tuple[str | None, str | None]:
    """(type name, why it could not be read).

    The reason is returned rather than swallowed for two reasons. python-pptx
    also returns None legitimately — for a graphic frame holding SmartArt — so
    collapsing both into None leaves a consumer unable to tell "no type" from
    "could not read the type". And it travels under the same `unreadable` key
    the per-shape guard below uses, so there is one marker for this class of gap
    rather than two for a consumer to reconcile.

    BaseShape.shape_type raises NotImplementedError for any subclass that
    doesn't override it, and Shape.shape_type raises for an sp that is neither
    placeholder, freeform, autoshape nor textbox. Letting that propagate cost
    the *whole* deck's style metadata over one odd shape — the caller catches
    per-file, not per-shape — so a criterion about branded text elsewhere on the
    slide would be graded with no evidence at all.
    """
    try:
        shape_type = shape.shape_type
    except Exception as e:
        return None, f"shape_type: {type(e).__name__}: {e}"
    return (str(shape_type) if shape_type else None), None


# Which master text style a placeholder falls back to when neither the shape,
# its layout placeholder nor the master placeholder states a size. PowerPoint
# keeps three on the master: titles, body, and everything else.
_TITLE_PH_TYPES = {"TITLE", "CENTER_TITLE", "VERTICAL_TITLE"}
_BODY_PH_TYPES = {
    "BODY",
    "SUBTITLE",
    "OBJECT",
    "CONTENT",
    "VERTICAL_BODY",
    "VERTICAL_OBJECT",
}


def _ph_category(ph_name: str) -> str:
    """Which master text style a placeholder inherits from: PowerPoint keeps
    three — titles, body, and everything else."""
    if ph_name in _TITLE_PH_TYPES:
        return "title"
    if ph_name in _BODY_PH_TYPES:
        return "body"
    return "other"


def _lvl_size_pt(element: Any, level: int) -> float | None:
    """Size declared for `level` in an element's a:lstStyle, in points.

    PowerPoint stores it per indent level as a:lvlNpPr/a:defRPr/@sz in
    hundredths of a point, so level 0 reads lvl1pPr.
    """
    if element is None:
        return None
    lst = element.find(f"{DML_NS}lstStyle")
    if lst is None:
        return None
    lvl = lst.find(f"{DML_NS}lvl{level + 1}pPr")
    if lvl is None:
        return None
    def_rpr = lvl.find(f"{DML_NS}defRPr")
    sz = def_rpr.get("sz") if def_rpr is not None else None
    return round(int(sz) / 100, 1) if sz else None


def _master_style_size_pt(master: Any, ph_type: str | None, level: int) -> float | None:
    """Size from the master's p:txStyles for this placeholder's category."""
    tx = master.element.find(f"{PML_NS}txStyles") if master is not None else None
    if tx is None:
        return None
    name = {"title": "titleStyle", "body": "bodyStyle"}.get(
        _ph_category(ph_type or ""), "otherStyle"
    )
    style = tx.find(f"{PML_NS}{name}")
    if style is None:
        return None
    lvl = style.find(f"{DML_NS}lvl{level + 1}pPr")
    def_rpr = lvl.find(f"{DML_NS}defRPr") if lvl is not None else None
    sz = def_rpr.get("sz") if def_rpr is not None else None
    return round(int(sz) / 100, 1) if sz else None


def _placeholder_para_size_pt(txbody: Any, level: int) -> float | None:
    """Size from a placeholder's own a:p/a:pPr/a:defRPr for this indent level.

    Matched on the paragraph's a:pPr@lvl (absent means level 0), so a
    multi-level body placeholder keeps each level's own size.
    """
    if txbody is None:
        return None
    for para in txbody.findall(f"{DML_NS}p"):
        ppr = para.find(f"{DML_NS}pPr")
        para_level = int((ppr.get("lvl") or "0") if ppr is not None else "0")
        if para_level != level:
            continue
        def_rpr = ppr.find(f"{DML_NS}defRPr") if ppr is not None else None
        sz = def_rpr.get("sz") if def_rpr is not None else None
        if sz:
            return round(int(sz) / 100, 1)
    return None


def _local_size_pt(
    txbody: Any, para: Any, list_style_source: str = "shape_list_style"
) -> tuple[float, str] | None:
    """The size declared on the paragraph itself or on its container's list
    style — the two links that need no placeholder chain, so a table cell can
    use them as well as a shape.

    a:pPr/a:defRPr applies to THIS paragraph, so it outranks the list style
    and everything below: a deck that overrides one paragraph was reported at
    the container's size instead of the size that renders.
    """
    ppr = para._p.find(f"{DML_NS}pPr")
    para_def = ppr.find(f"{DML_NS}defRPr") if ppr is not None else None
    para_sz = para_def.get("sz") if para_def is not None else None
    if para_sz:
        return round(int(para_sz) / 100, 1), "paragraph_default"

    own = _lvl_size_pt(txbody, para.level)
    return (own, list_style_source) if own is not None else None


def _presentation_default_size_pt(master: Any, level: int) -> tuple[float, str] | None:
    """Size from p:presentation/p:defaultTextStyle — the last link in the chain.

    Deliberately last: it is the weakest declaration, and a shape's own list
    style outranks it. Reaching it earlier reported 18pt for text that renders
    at 28. Reaching it at all matters for a plain text box that states nothing
    locally, whose size was otherwise left "inherited".
    """
    try:
        prs = master.part.package.presentation_part._element
        default = prs.find(f"{PML_NS}defaultTextStyle")
    except Exception:
        return None
    size = _lvl_size_pt_from(default, level)
    return (size, "presentation_default") if size is not None else None


def _lvl_size_pt_from(style_el: Any, level: int) -> float | None:
    """Size for `level` from an element that holds a:lvlNpPr children."""
    if style_el is None:
        return None
    lvl = style_el.find(f"{DML_NS}lvl{level + 1}pPr")
    def_rpr = lvl.find(f"{DML_NS}defRPr") if lvl is not None else None
    sz = def_rpr.get("sz") if def_rpr is not None else None
    return round(int(sz) / 100, 1) if sz else None


def _inherited_size_pt(shape: Any, para: Any, layout: Any) -> tuple[float, str] | None:
    """The size a run with no explicit size actually renders at, and its source.

    Walks the chain PowerPoint does, and the ORDER matters: the paragraph's own
    a:pPr/a:defRPr, then the shape's own a:lstStyle, then the layout
    placeholder, the master placeholder, the master txStyles, and last the
    presentation default. A placeholder copied onto a slide keeps its list
    style while losing placeholder status, so reaching the presentation default
    earlier would report 18pt for text that renders at 28.

    Returns (size_pt, source), or None when no link declares one.
    """
    try:
        level = para.level

        local = _local_size_pt(shape._element.find(f"{PML_NS}txBody"), para)
        if local is not None:
            return local

        if not shape.is_placeholder or layout is None:
            # A plain text box has no placeholder chain, but the presentation
            # still declares a default text size that it renders at.
            master = getattr(layout, "slide_master", None) if layout else None
            return _presentation_default_size_pt(master, para.level)
        idx = shape.placeholder_format.idx
        ph_type = shape.placeholder_format.type
        ph_name = getattr(ph_type, "name", None) or str(ph_type or "")

        for lp in layout.placeholders:
            if lp.placeholder_format.idx != idx:
                continue
            lp_body = lp._element.find(f"{PML_NS}txBody")
            # The layout placeholder can declare the size on its own
            # paragraphs rather than its list style, and that is the more
            # specific of the two — reading only the list style fell through
            # to the master and reported a size the layout overrides.
            size = _placeholder_para_size_pt(lp_body, level)
            if size is None:
                size = _lvl_size_pt(lp_body, level)
            if size is not None:
                return size, "layout_placeholder"
            break

        master = getattr(layout, "slide_master", None)
        if master is None:
            return None
        # Exact type first, then category. Category alone is needed because a
        # slide's content placeholder is OBJECT while the master's is BODY, but
        # it is too coarse on its own: date, footer and slide-number are all
        # "other", so the first of them would answer for all three.
        candidates = [
            mp
            for mp in master.placeholders
            if (getattr(mp.placeholder_format.type, "name", None) or "") == ph_name
        ] or [
            mp
            for mp in master.placeholders
            if _ph_category(getattr(mp.placeholder_format.type, "name", None) or "")
            == _ph_category(ph_name)
        ]
        for mp in candidates:
            mp_body = mp._element.find(f"{PML_NS}txBody")
            # Paragraph default before list style, as for the layout above.
            size = _placeholder_para_size_pt(mp_body, level)
            if size is None:
                size = _lvl_size_pt(mp_body, level)
            if size is not None:
                return size, "master_placeholder"
            break

        size = _master_style_size_pt(master, ph_name, level)
        if size is not None:
            return size, "master_text_styles"
        return _presentation_default_size_pt(master, level)
    except Exception as e:
        # Inheritance is best-effort: a deck with an odd placeholder graph must
        # fall back to "inherited" rather than fail the whole extraction. The
        # fallback is indistinguishable from a genuine "nothing declares a
        # size", so it is logged — otherwise a resolver that breaks on a whole
        # family of decks looks exactly like decks that state no size.
        logger.warning(
            f"[TRANSFORM] pptx size inheritance failed: {type(e).__name__}: {e}"
        )
        return None


def _extract_shape_metadata(shape: Any, layout: Any = None) -> dict[str, Any]:
    """Extract metadata from a single shape."""
    shape_type, type_unreadable = _shape_type_name(shape)
    result: dict[str, Any] = {
        "shape_name": shape.name,
        "shape_type": shape_type,
        "left": _emu_to_pt(shape.left),
        "top": _emu_to_pt(shape.top),
        "width": _emu_to_pt(shape.width),
        "height": _emu_to_pt(shape.height),
    }
    if type_unreadable:
        # Same key the per-shape guard below uses, so a consumer has one marker
        # for this class of gap rather than two to reconcile.
        result["unreadable"] = type_unreadable

    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            runs = [_extract_run_style(run) for run in para.runs]
            # Fill in the size the run actually renders at when it declares
            # none of its own, and say where the value came from so a reader
            # can tell an inherited size from one set on the run.
            inherited = None
            if any(r["font_size_pt"] is None for r in runs):
                inherited = _inherited_size_pt(shape, para, layout)
            if inherited is not None:
                size, source = inherited
                for r in runs:
                    if r["font_size_pt"] is None:
                        r["font_size_pt"] = size
                        r["font_size_source"] = source
            paragraphs.append(
                {
                    "alignment": str(para.alignment) if para.alignment else None,
                    "level": para.level,
                    "text_runs": runs,
                }
            )
        result["paragraphs"] = paragraphs

    # A table is a GraphicFrame, so has_text_frame is False and every cell's
    # text and styling was skipped entirely — a criterion like "the table
    # content is 28pt" had nothing to read even though the cells state it
    # explicitly. Cell paragraphs are appended as the shape's own so the
    # per-slide style summary counts them like any other run.
    if getattr(shape, "has_table", False):
        table_paragraphs: list[dict[str, Any]] = []
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        runs = [_extract_run_style(run) for run in para.runs]
                        if not runs:
                            continue
                        # A cell has no placeholder chain, but it can still
                        # carry a paragraph default or its own list style,
                        # and those are what render.
                        if any(r["font_size_pt"] is None for r in runs):
                            local = _local_size_pt(
                                cell._tc.find(f"{DML_NS}txBody"),
                                para,
                                "cell_list_style",
                            )
                            if local is None and layout is not None:
                                # A cell states nothing of its own, so it falls
                                # to the presentation default like any other
                                # text. Stopping here reported "inherited" for
                                # table text while identical text in a plain
                                # box resolved.
                                local = _presentation_default_size_pt(
                                    getattr(layout, "slide_master", None), para.level
                                )
                            if local is not None:
                                size, source = local
                                for r in runs:
                                    if r["font_size_pt"] is None:
                                        r["font_size_pt"] = size
                                        r["font_size_source"] = source
                        table_paragraphs.append(
                            {
                                "alignment": str(para.alignment)
                                if para.alignment
                                else None,
                                "level": para.level,
                                "in_table": True,
                                "text_runs": runs,
                            }
                        )
        except Exception as e:
            # Same rule as the shape-type guard: one odd table must not cost
            # the whole deck's metadata.
            result["unreadable"] = f"table: {type(e).__name__}: {e}"
        if table_paragraphs:
            result["paragraphs"] = result.get("paragraphs", []) + table_paragraphs

    return result


# Guard against a malformed deck with a cyclic or absurdly deep group nesting.
_MAX_GROUP_DEPTH = 8


def _flatten_shapes(shapes: Any, depth: int = 0) -> list[Any]:
    """Yield shapes, descending into group shapes.

    slide.shapes lists a GroupShape but not its children, so grouped text is
    invisible to a plain iteration. Opt-in only, so the older verifiers keep
    the shape list they were written against — see pptx_to_style_metadata.
    """
    flat: list[Any] = []
    for shape in shapes:
        flat.append(shape)
        if depth >= _MAX_GROUP_DEPTH:
            continue
        # shape_type can raise on exotic shapes, so key off the child
        # collection that GroupShape actually exposes.
        children = getattr(shape, "shapes", None)
        if children is not None:
            flat.extend(_flatten_shapes(children, depth + 1))
    return flat


def pptx_to_style_metadata(
    file_bytes: bytes, file_name: str, include_grouped_shapes: bool = False
) -> dict[str, Any]:
    """Extract style metadata from a PPTX file.

    Args:
        file_bytes: Raw .pptx bytes.
        file_name: Filename, included in the output for context.
        include_grouped_shapes: Descend into group shapes. Defaults off because
            two pre-existing evals (pptx_style_verifier and
            pptx_style_verifier_apex_v2) json.dumps this dict wholesale, and
            adding a slide's group children plus the group shapes themselves
            changes any criterion that counts shapes — a 3-shape slide becomes
            5. Only the generic multi-representation judge opts in, where the
            output is a collapsed style summary that no criterion counts.

    Returns:
        {
            "file_name": str,
            "slide_width_pt": float,
            "slide_height_pt": float,
            "slide_count": int,
            "slides": [
                {
                    "index": int,
                    "title": str | None,
                    "layout_name": str,
                    "shapes": [...]
                }
            ]
        }
    """
    prs = Presentation(BytesIO(file_bytes))

    slide_width_pt = _emu_to_pt(prs.slide_width)
    slide_height_pt = _emu_to_pt(prs.slide_height)

    slides_meta: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        # Extract title from the title placeholder if present
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text

        layout_name = slide.slide_layout.name if slide.slide_layout else None

        # Degrade per shape, not per file. python-pptx raises on several
        # properties of shapes it cannot model, and the caller only catches at
        # file granularity — so one odd shape would otherwise discard every
        # other shape's fonts and colours along with it. The unreadable shape is
        # recorded rather than dropped so a consumer can see the gap.
        raw_shapes = (
            _flatten_shapes(slide.shapes)
            if include_grouped_shapes
            else list(slide.shapes)
        )
        shapes: list[dict[str, Any]] = []
        for shape in raw_shapes:
            try:
                shapes.append(_extract_shape_metadata(shape, slide.slide_layout))
            except Exception as e:
                logger.warning(
                    f"[TRANSFORM] unreadable shape in {file_name}: "
                    f"{type(e).__name__}: {e}"
                )
                shapes.append(
                    {
                        "shape_name": None,
                        "shape_type": None,
                        "unreadable": f"{type(e).__name__}: {e}",
                    }
                )

        slides_meta.append(
            {
                "index": idx,
                "title": title,
                "layout_name": layout_name,
                "shapes": shapes,
            }
        )

    return {
        "file_name": file_name,
        "slide_width_pt": slide_width_pt,
        "slide_height_pt": slide_height_pt,
        "slide_count": len(slides_meta),
        "slides": slides_meta,
    }


# ==========================================================================
# TRANSFORMATION WRAPPER (for the generic multi-representation judge)
# ==========================================================================
# The dict above is consumed directly by pptx_style_verifier, which json.dumps
# it wholesale. That verifier grades one deck per call, so size doesn't matter
# there. The generic judge is different: its style metadata shares a capped
# token budget with everything else in the prompt, and on a spreadsheet an
# unbounded dump cost 45MB and lost 7 of 8 sheets to truncation. So this emits
# a compact per-slide style summary instead — runs collapsed by identical
# style, which is what font/size/colour criteria actually need.

# Distinct style combinations listed per slide before summarising the rest.
_MAX_STYLES_PER_SLIDE = 12

# Characters of sample text kept per style, purely to help the judge tie a
# style to what it applies to.
_SAMPLE_TEXT_CHARS = 60

# Shape names listed per style group. A sample, and said to be one when it is
# truncated — see the shapes_listed attribute.
_MAX_SHAPE_NAMES = 3


def _run_style_key(run: dict[str, Any]) -> tuple[Any, ...]:
    return (
        run.get("font_name"),
        run.get("font_size_pt"),
        run.get("font_color_rgb") or run.get("font_color_theme"),
        run.get("bold"),
        run.get("italic"),
        # Part of the key, not just the output: a size the run states and one
        # resolved from a layout are different facts, and collapsing them
        # would report a single group whose provenance is true of only some
        # of its runs.
        run.get("font_size_source"),
    )


def _style_attrs(key: tuple[Any, ...]) -> str:
    font_name, size, color, bold, italic, size_source = key
    attrs: list[str] = []
    # None means "inherited from the slide layout/master" rather than unset,
    # so say that explicitly instead of silently omitting it — otherwise the
    # judge can't tell "not bold" from "not specified here".
    attrs.append(
        f'font_name="{xml_escape(str(font_name))}"'
        if font_name
        else 'font_name="inherited"'
    )
    attrs.append(
        f'font_size_pt="{size}"' if size is not None else 'font_size_pt="inherited"'
    )
    # Where the size came from, so a reader can tell a size the run states
    # from one this extractor walked the inheritance chain for.
    if size_source:
        attrs.append(f'font_size_pt_source="{xml_escape(str(size_source))}"')
    if color:
        # A theme slot is prefixed so the judge cannot mistake it for a hex.
        value = str(color)
        if not value.startswith("#"):
            value = f"theme:{value}"
        attrs.append(f'font_color="{xml_escape(value)}"')
    # Same reasoning as the fonts above, which this previously failed to
    # follow: omitting bold/italic when None left the judge unable to tell
    # inherited-bold from not-bold, on criteria whose subject is exactly that.
    attrs.append(
        f'bold="{str(bool(bold)).lower()}"' if bold is not None else 'bold="inherited"'
    )
    attrs.append(
        f'italic="{str(bool(italic)).lower()}"'
        if italic is not None
        else 'italic="inherited"'
    )
    return " ".join(attrs)


def _pptx_style_summary_xml(data: dict[str, Any]) -> str:
    lines: list[str] = [
        f'<presentation slides="{data.get("slide_count", 0)}" '
        f'slide_width_pt="{data.get("slide_width_pt")}" '
        f'slide_height_pt="{data.get("slide_height_pt")}">'
    ]

    for slide in data.get("slides", []):
        title = slide.get("title")
        attrs = f'index="{slide.get("index", 0) + 1}"'
        if slide.get("layout_name"):
            attrs += f' layout="{xml_escape(str(slide["layout_name"]))}"'
        if title:
            attrs += f' title="{xml_escape(title[:_SAMPLE_TEXT_CHARS])}"'
        lines.append(f"  <slide {attrs}>")

        # Collapse every run on the slide by identical style signature. A
        # deck reuses a handful of styles, so this stays tiny regardless of
        # how much text the slide holds.
        counts: Counter[tuple[Any, ...]] = Counter()
        samples: dict[tuple[Any, ...], str] = {}
        shapes_for: dict[tuple[Any, ...], set[str]] = {}
        for shape in slide.get("shapes", []):
            for para in shape.get("paragraphs", []):
                for run in para.get("text_runs", []):
                    key = _run_style_key(run)
                    counts[key] += 1
                    text = (run.get("text") or "").strip()
                    if text and key not in samples:
                        samples[key] = text[:_SAMPLE_TEXT_CHARS]
                    if shape.get("shape_name"):
                        shapes_for.setdefault(key, set()).add(shape["shape_name"])

        # Two gaps, reported separately. A shape that failed entirely has no
        # runs below, so its absence needs disclosing; a shape whose TYPE alone
        # could not be read still has its runs listed, so calling those
        # unrepresented would tell the judge to discount real evidence.
        def _represented(shape: dict[str, Any]) -> bool:
            return any(para.get("text_runs") for para in shape.get("paragraphs", []))

        flagged = [s for s in slide.get("shapes", []) if s.get("unreadable")]
        unreadable = [str(s.get("unreadable")) for s in flagged if not _represented(s)]
        partial = [str(s.get("unreadable")) for s in flagged if _represented(s)]

        if unreadable:
            reason = xml_escape(unreadable[0][:_SAMPLE_TEXT_CHARS])
            lines.append(
                f'    <unreadable_shapes count="{len(unreadable)}" '
                f'first_error="{reason}">Text and styles of these shapes are '
                f"NOT represented below; draw no conclusion about them from "
                f"their absence.</unreadable_shapes>"
            )
        if partial:
            reason = xml_escape(partial[0][:_SAMPLE_TEXT_CHARS])
            lines.append(
                f'    <partial_shapes count="{len(partial)}" '
                f'first_error="{reason}">The SHAPE TYPE of these shapes could '
                f"not be read, so draw no conclusion about what kind of object "
                f"they are. Their text and text styles ARE listed below and can "
                f"be relied on.</partial_shapes>"
            )

        for key, count in counts.most_common(_MAX_STYLES_PER_SLIDE):
            attrs = _style_attrs(key) + f' runs="{count}"'
            all_shapes = sorted(shapes_for.get(key, ()))
            shapes = all_shapes[:_MAX_SHAPE_NAMES]
            if shapes:
                attrs += f' shapes="{xml_escape(", ".join(shapes))}"'
                # Say the list is a sample. Unqualified, a criterion about
                # "the title and the footer" could read a 3-name list as the
                # complete set of shapes carrying this style.
                if len(all_shapes) > len(shapes):
                    attrs += f' shapes_listed="{len(shapes)} of {len(all_shapes)}"'
            sample = samples.get(key)
            if sample:
                lines.append(
                    f"    <text_style {attrs}>{xml_escape(sample)}</text_style>"
                )
            else:
                lines.append(f"    <text_style {attrs} />")

        omitted = len(counts) - min(len(counts), _MAX_STYLES_PER_SLIDE)
        if omitted > 0:
            lines.append(f"    <!-- {omitted} rarer style combinations omitted -->")

        lines.append("  </slide>")

    lines.append("</presentation>")
    return "\n".join(lines)


async def _extract_pptx_style_text(file_bytes: bytes, file_name: str) -> str:
    # Parsing is CPU-bound; keep it off the shared grading event loop.
    data = await asyncio.to_thread(pptx_to_style_metadata, file_bytes, file_name, True)
    return _pptx_style_summary_xml(data)


async def pptx_to_style_metadata_output(
    file_bytes: bytes, file_name: str
) -> TransformationOutput:
    """Per-slide font/size/colour summary as judge-readable XML.

    Answers criteria like "titles are 30pt on slides 1 and 7" or "the deck
    uses Arial Black" without the judge having to read font sizes off a
    rendered slide image.
    """
    if not is_ooxml_package(file_bytes):
        logger.info(
            f"[TRANSFORM] {file_name} is not an OOXML package — no pptx style "
            f"metadata available"
        )
        return TransformationOutput(
            text=(
                '<presentation unsupported="true" data-degraded="true">Style metadata requires a '
                ".pptx (OOXML) file; legacy .ppt and OpenDocument .odp are not "
                "supported by this extractor.</presentation>\n"
            )
        )

    text = await cached_style_text(
        file_bytes, file_name, "pptx", _extract_pptx_style_text
    )
    return TransformationOutput(text=text)
