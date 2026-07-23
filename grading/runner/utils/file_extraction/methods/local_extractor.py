"""
Local file extractor using Python libraries for fast, lightweight extraction.

This is the offline last-resort tier: it runs when the preferred cloud
extractors (Mercor's document cache, then Reducto) are unavailable or fail.
Fast and dependency-light, but text-only — no OCR or layout.

Supported formats:
- XLSX: openpyxl (recalculates formulas via LibreOffice when no cached value
  is present, with optional chart extraction via LibreOffice PDF conversion)
- PPTX: python-pptx
- DOCX: python-docx
- PDF: pypdf (text-only offline fallback; no OCR/layout — used when
  Reducto/Mercor aren't configured)
- CSV: built-in csv module (always available)
- TXT: built-in (always available)
"""

import csv
import io
import math
import re
import shutil
import zipfile
from decimal import ROUND_HALF_UP, Decimal, InvalidOperation
from importlib.util import find_spec
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

import openpyxl
import xlrd
from docx import Document
from docx.oxml.ns import qn
from loguru import logger
from pptx import Presentation
from pypdf import PdfReader
from tenacity import retry, stop_after_attempt, wait_fixed

from ..base import BaseFileExtractor
from ..constants import SPREADSHEET_EXTENSIONS
from ..types import ExtractedContent, ImageMetadata, SubArtifact
from ..utils.chart_extraction import (
    convert_xlsx_to_pdf,
    evaluate_excel_formulas_with_libreoffice,
    find_libreoffice,
    has_charts_in_xlsx,
    pdf_to_base64_images,
)


def _build_comment_anchor_map(doc: Any) -> dict[str, dict[str, Any]]:
    """Map comment IDs to the paragraph index and anchored text they bracket.

    DOCX stores comments in two files: ``word/comments.xml`` holds the comment
    bodies, while ``word/document.xml`` carries ``<w:commentRangeStart>`` /
    ``<w:commentRangeEnd>`` markers around the anchored run(s). python-docx's
    high-level paragraph iteration strips those markers, which leaves the LLM
    judge unable to grade ``"comment attached to paragraph 17"`` style criteria.

    Walks each paragraph's underlying lxml tree once, tracks open ranges across
    paragraph boundaries, and returns ``{comment_id: {"paragraph_index",
    "anchored_text"}}``. Cross-paragraph ranges are recorded against the
    paragraph where they STARTED.
    """
    anchor_map: dict[str, dict[str, Any]] = {}
    open_ranges: dict[str, dict[str, Any]] = {}
    range_start_tag = qn("w:commentRangeStart")
    range_end_tag = qn("w:commentRangeEnd")
    text_tag = qn("w:t")
    id_attr = qn("w:id")

    for para_idx, para in enumerate(doc.paragraphs, start=1):
        for elem in para._element.iter():  # noqa: SLF001 — lxml access is the only stable path
            tag = elem.tag
            if tag == range_start_tag:
                cid = elem.get(id_attr)
                if cid is not None:
                    open_ranges[cid] = {
                        "paragraph_index": para_idx,
                        "text_parts": [],
                    }
            elif tag == range_end_tag:
                cid = elem.get(id_attr)
                if cid is not None and cid in open_ranges:
                    entry = open_ranges.pop(cid)
                    anchor_map[cid] = {
                        "paragraph_index": entry["paragraph_index"],
                        "anchored_text": "".join(entry["text_parts"]).strip(),
                    }
            elif tag == text_tag and elem.text:
                for entry in open_ranges.values():
                    entry["text_parts"].append(elem.text)

    for cid, entry in open_ranges.items():
        anchor_map[cid] = {
            "paragraph_index": entry["paragraph_index"],
            "anchored_text": "".join(entry["text_parts"]).strip() or "(unclosed range)",
        }

    return anchor_map


def _xlsx_has_uncached_formulas(source: Path | bytes, source_name: str = "") -> bool:
    """True iff a formula cell reads None under data_only (no cached result).

    openpyxl has no formula engine, so formula cells a real spreadsheet never
    saved read None and vanish from the text. Correlating the data_only/formula
    loads by coordinate avoids the expensive LibreOffice recalc on already-cached
    workbooks. Streams from a Path (no full read into memory) when given one.
    """

    def _open(**kw: Any) -> Any:
        src = io.BytesIO(source) if isinstance(source, bytes) else source
        return openpyxl.load_workbook(src, **kw)

    try:
        # Formula coordinates first — usually few, even on a bloated used range —
        # so a sheet full of blank cells never materializes into a giant set.
        # read_only streams both loads (data_type=='f' is still reported).
        wb_formulas = _open(data_only=False, read_only=True)
        try:
            formula_coords = {
                (sheet_name, cell.coordinate)
                for sheet_name in wb_formulas.sheetnames
                for row in wb_formulas[sheet_name].iter_rows()
                for cell in row
                if cell.data_type == "f"
                or (cell.value is not None and str(cell.value).startswith("="))
            }
        finally:
            wb_formulas.close()

        if not formula_coords:
            return False

        # A formula cell reading None under data_only had no saved result.
        wb_values = _open(data_only=True, read_only=True)
        try:
            return any(
                cell.value is None
                and (sheet_name, getattr(cell, "coordinate", None)) in formula_coords
                for sheet_name in wb_values.sheetnames
                for row in wb_values[sheet_name].iter_rows()
                for cell in row
            )
        finally:
            wb_values.close()
    except Exception as e:
        # On uncertainty assume uncached (True): recalc then fixes or flags the
        # values, whereas False would drop real formulas with no signal.
        logger.warning(
            f"[LOCAL] {source_name or 'workbook'}: could not determine "
            f"uncached-formula status ({type(e).__name__}: {e}); assuming "
            f"uncached and recalculating."
        )
        return True


# Excel locale tag [$SYM-hexlocale]: the '$' after '[' is a marker, not a symbol;
# SYM (between '[$' and '-') is empty for a bare locale like [$-409] (not currency).
_CURRENCY_LOCALE_TAG_RE = re.compile(r"\[\$([^\]-]*)-[0-9A-Fa-f]+\]")
# Any bracket tag ([$...], [Red], [>0]); stripped before the bare-symbol search.
_BRACKET_TAG_RE = re.compile(r"\[[^\]]*\]")
_NUMBER_FORMAT_CURRENCY_RE = re.compile(r"[$€£¥]")
# A cell worth surfacing in the Number Formats block: currency/percent/grouped
# (skips dates, text, General). Gates the annotation; rendering is separate.
_NUMERIC_DISPLAY_FORMAT_RE = re.compile(r"[$€£¥%]|#,##0")
# Shapes we render EXACTLY (post tag/symbol strip): optional #,## grouping, digits,
# optional '0' decimals, optional %. Anything else → None (not reproduced).
_RENDERABLE_NUMBER_RE = re.compile(r"(?:#,##)?[0#]+(?:\.0+)?%?")
# Accounting fill/alignment tokens (_x reserves a char's width, *x repeats a fill
# char) are cosmetic column padding — strip them so accounting formats render.
_ALIGN_FILL_RE = re.compile(r"[_*].")
# Currency symbols arrive escaped (\$) after recalc or quoted ("$") from Excel;
# both display the bare symbol, so normalize to $ before matching.
_ESCAPED_CURRENCY_RE = re.compile(r"\\([$€£¥])")
_QUOTED_CURRENCY_RE = re.compile(r'"([$€£¥])"')
# A conditional section ([>=1000]…) makes Excel choose a section by value; we
# don't evaluate that, so a format carrying one isn't safely renderable.
_CONDITION_TAG_RE = re.compile(r"\[[<>=]")


def _is_numeric_display_format(number_format: str) -> bool:
    """True for currency/percent/grouped formats (not dates/text). The $ in a bare
    locale tag ([$-409]) is a marker, not currency, so strip tags first.
    """
    locale_match = _CURRENCY_LOCALE_TAG_RE.search(number_format)
    if locale_match and locale_match.group(1):
        return True
    return bool(
        _NUMERIC_DISPLAY_FORMAT_RE.search(_BRACKET_TAG_RE.sub("", number_format))
    )


def _render_number_format(value: Any, number_format: str | None) -> str | None:
    """Render a numeric cell per its Excel format, or None if we can't do it
    exactly. None is safe: the raw value + [fmt: …] tag always stand alongside.
    """
    if not number_format or number_format == "General":
        return None
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        return None

    # Negatives/zeros with their OWN section (e.g. accounting parens) display from
    # it, not the positive one — we don't reproduce that literal wrapping.
    number_format = _ESCAPED_CURRENCY_RE.sub(r"\1", number_format)
    number_format = _QUOTED_CURRENCY_RE.sub(r"\1", number_format)
    if _CONDITION_TAG_RE.search(number_format):
        return None
    sections = number_format.split(";")
    if value < 0 and len(sections) > 1 and sections[1].strip():
        return None
    if value == 0 and len(sections) > 2 and sections[2].strip():
        return None

    positive = _BRACKET_TAG_RE.sub("", sections[0])
    locale_match = _CURRENCY_LOCALE_TAG_RE.search(sections[0])
    symbol_match = _NUMBER_FORMAT_CURRENCY_RE.search(positive)
    symbol = (locale_match.group(1) if locale_match else "") or (
        symbol_match.group(0) if symbol_match else ""
    )
    core = positive.replace(symbol, "", 1) if symbol else positive
    core = _ALIGN_FILL_RE.sub("", core).strip()

    # Only shapes we reproduce EXACTLY; anything else → None (raw value + [fmt] tag
    # still stand). See _RENDERABLE_NUMBER_RE for what's rejected.
    if not _RENDERABLE_NUMBER_RE.fullmatch(core):
        return None

    decimals = len(core.split(".", 1)[1].rstrip("%")) if "." in core else 0
    is_percent = core.endswith("%")
    try:
        magnitude = Decimal(str(abs(value))) * (100 if is_percent else 1)
        quantized = magnitude.quantize(Decimal(1).scaleb(-decimals), ROUND_HALF_UP)
    except (InvalidOperation, ValueError):
        # Out-of-range value (quantize needs >28 sig digits) — an optional
        # annotation must never abort extraction, so leave the raw value to stand.
        return None
    # Group only when the format asks for it (#,##); $0.00 / 0% display ungrouped.
    grouped = "," if "#,##" in core else ""
    number = f"{quantized:{grouped}.{decimals}f}"
    sign = "-" if value < 0 else ""
    return f"{sign}{symbol}{number}{'%' if is_percent else ''}"


class LocalExtractor(BaseFileExtractor):
    """
    Local extractor for quick content extraction using Python libraries.

    This extractor is fast but provides basic text extraction. It's designed
    for change detection rather than high-quality content extraction.
    """

    def __init__(self):
        """Initialize the local extractor"""
        self._supported_extensions = set()

        # Check for openpyxl
        if find_spec("openpyxl") is not None:
            self._has_openpyxl = True
            self._supported_extensions.update(SPREADSHEET_EXTENSIONS)
            logger.debug("LocalExtractor: openpyxl available for XLSX files")
        else:
            self._has_openpyxl = False
            logger.debug("LocalExtractor: openpyxl not available")

        # Check for python-pptx
        if find_spec("pptx") is not None:
            self._has_pptx = True
            self._supported_extensions.update([".pptx"])
            logger.debug("LocalExtractor: python-pptx available for PPTX files")
        else:
            self._has_pptx = False
            logger.debug("LocalExtractor: python-pptx not available")

        # Check for python-docx
        if find_spec("docx") is not None:
            self._has_docx = True
            self._supported_extensions.update([".docx"])
            logger.debug("LocalExtractor: python-docx available for DOCX files")
        else:
            self._has_docx = False
            logger.debug("LocalExtractor: python-docx not available")

        # Check for xlrd (for .xls files)
        if find_spec("xlrd") is not None:
            self._has_xlrd = True
            self._supported_extensions.add(".xls")
            logger.debug("LocalExtractor: xlrd available for XLS files")
        else:
            self._has_xlrd = False
            logger.debug("LocalExtractor: xlrd not available")

        # Check for csv (built-in, always available)
        if find_spec("csv") is not None:
            self._has_csv = True
            self._supported_extensions.update([".csv"])
            logger.debug("LocalExtractor: csv available for CSV files")
        else:
            self._has_csv = False
            logger.debug("LocalExtractor: csv not available")

        # Offline PDF fallback (text-only, no OCR): keeps PDFs readable when
        # Reducto/Mercor are unconfigured or fail.
        if find_spec("pypdf") is not None:
            self._has_pypdf = True
            self._supported_extensions.add(".pdf")
            logger.debug("LocalExtractor: pypdf available for PDF files")
        else:
            self._has_pypdf = False
            logger.debug("LocalExtractor: pypdf not available")

    @property
    def name(self) -> str:
        return "local_python_libs"

    def supports_file_type(self, file_extension: str) -> bool:
        """Check if this extractor supports the given file type"""
        return file_extension.lower() in self._supported_extensions

    async def extract_from_file(
        self,
        file_path: Path,
        *,
        include_images: bool = True,
        sub_artifact_index: int | None = None,
    ) -> ExtractedContent:
        """
        Extract content from a file using local Python libraries.

        This provides basic text extraction for change detection.
        """
        file_ext = file_path.suffix.lower()

        # Route .xls files to xlrd extractor (openpyxl doesn't support .xls)
        if file_ext == ".xls" and self._has_xlrd:
            return await self._extract_xls(file_path, sub_artifact_index)
        elif file_ext in SPREADSHEET_EXTENSIONS and self._has_openpyxl:
            return await self._extract_xlsx(file_path, sub_artifact_index)
        elif file_ext == ".pptx" and self._has_pptx:
            return await self._extract_pptx(file_path, sub_artifact_index)
        elif file_ext == ".docx" and self._has_docx:
            return await self._extract_docx(file_path, sub_artifact_index)
        elif file_ext == ".pdf" and self._has_pypdf:
            return await self._extract_pdf(file_path, sub_artifact_index)
        elif file_ext == ".csv" and self._has_csv:
            return await self._extract_csv(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

    def _get_hidden_sheets_from_xlsx(self, file_path: Path) -> set[str]:
        """Extract hidden sheet names from xlsx by parsing workbook.xml directly.

        This is lightweight and works regardless of read_only mode, since it
        only parses the workbook metadata XML, not the cell data.
        """
        hidden_sheets: set[str] = set()
        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                with zf.open("xl/workbook.xml") as f:
                    tree = ET.parse(f)
                    root = tree.getroot()

                    # xlsx uses Office Open XML namespace
                    ns = {
                        "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                    }

                    for sheet in root.findall(".//main:sheet", ns):
                        state = sheet.get("state", "visible")
                        if state in ("hidden", "veryHidden"):
                            sheet_name = sheet.get("name")
                            if sheet_name:
                                hidden_sheets.add(sheet_name)
                                logger.debug(
                                    f"Detected hidden sheet '{sheet_name}' (state: {state})"
                                )
        except Exception as e:
            logger.warning(f"Failed to parse workbook.xml for hidden sheets: {e}")

        return hidden_sheets

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_fixed(1),
        reraise=True,
    )
    async def _extract_xlsx(
        self, file_path: Path, sub_artifact_index: int | None = None
    ) -> ExtractedContent:
        """Extract content from XLSX file using openpyxl, with optional chart extraction."""
        if openpyxl is None:
            raise ImportError("openpyxl is required for XLSX extraction")

        try:
            # Get hidden sheets by parsing workbook.xml directly
            hidden_sheets = self._get_hidden_sheets_from_xlsx(file_path)

            # Uncached formula cells read as None under data_only and vanish from
            # the text; recalc via LibreOffice first so they stay visible.
            workbook_source: Path | io.BytesIO = file_path
            formulas_unresolved = False
            if _xlsx_has_uncached_formulas(file_path, file_path.name):
                # Read bytes only when recalc fires (cached case streams). Pass the
                # real suffix (.xlsm) or the recalc temp file defaults to .xlsx.
                recalculated_bytes = await evaluate_excel_formulas_with_libreoffice(
                    file_path.read_bytes(), suffix=file_path.suffix
                )
                if recalculated_bytes:
                    workbook_source = io.BytesIO(recalculated_bytes)
                else:
                    formulas_unresolved = True
                    logger.warning(
                        f"[LOCAL] {file_path.name} has formula cells with no "
                        "cached value and LibreOffice recalculation is "
                        "unavailable — extracting cached values only "
                        "(formula cells may appear blank)."
                    )

            # When recalc was unavailable, note it in the text so blank formula
            # cells aren't graded as the agent's own empty output.
            unresolved_banner = (
                "=== Extraction Note ===\n[formula cells could not be "
                "recalculated — LibreOffice unavailable; blanks may be "
                "uncomputed formulas, not missing values]\n\n"
                if formulas_unresolved
                else ""
            )

            # read_only streams the workbook instead of loading it all into memory
            wb = openpyxl.load_workbook(workbook_source, data_only=True, read_only=True)

            # Collect authored formula text keyed by (sheet, coordinate) so formula
            # criteria stay verifiable. Scoped to rendered sheets; coord joins the value pass.
            if sub_artifact_index is not None:
                _wanted_sheets = (
                    {wb.sheetnames[sub_artifact_index]}
                    if 0 <= sub_artifact_index < len(wb.sheetnames)
                    else set()
                )
            else:
                _wanted_sheets = set(wb.sheetnames)
            _wanted_sheets -= hidden_sheets

            formula_text: dict[tuple[str, str], str] = {}
            try:
                wb_formulas = openpyxl.load_workbook(
                    file_path, data_only=False, read_only=True
                )
                try:
                    for _sheet_name in wb_formulas.sheetnames:
                        if _sheet_name not in _wanted_sheets:
                            continue
                        for _row in wb_formulas[_sheet_name].iter_rows():
                            for _cell in _row:
                                # Select formula cells by data_type 'f'; the "=" check
                                # below is a sanity guard (a real formula starts with "=").
                                if _cell.data_type != "f":
                                    continue
                                _v = _cell.value
                                _f = (
                                    _v
                                    if isinstance(_v, str)
                                    else getattr(_v, "text", None)
                                )
                                _coord = getattr(_cell, "coordinate", None)
                                if (
                                    _coord
                                    and isinstance(_f, str)
                                    and _f.startswith("=")
                                ):
                                    formula_text[(_sheet_name, _coord)] = _f
                finally:
                    wb_formulas.close()
            except Exception:
                logger.opt(exception=True).warning(
                    f"[LOCAL] {file_path.name}: formula-text pass failed; "
                    "rendering cached values only."
                )
                formula_text = {}

            sub_artifacts = []
            full_text_parts = []
            skipped_hidden_sheet: str | None = (
                None  # Track if we skipped requested sheet
            )

            try:
                for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                    # If specific sub-artifact requested, skip others
                    if (
                        sub_artifact_index is not None
                        and sheet_idx != sub_artifact_index
                    ):
                        continue

                    # Skip hidden sheets
                    if sheet_name in hidden_sheets:
                        logger.debug(f"Skipping hidden sheet '{sheet_name}'")
                        # Track if this was the specifically requested sheet
                        if (
                            sub_artifact_index is not None
                            and sheet_idx == sub_artifact_index
                        ):
                            skipped_hidden_sheet = sheet_name
                        continue

                    sheet = wb[sheet_name]

                    # Full Cell objects (not values_only) so number formats and authored
                    # formulas emit as separate blocks, leaving the cell grid unchanged.
                    sheet_text_lines = []
                    format_lines = []
                    formula_lines = []
                    for row in sheet.iter_rows():
                        row_values = []
                        for cell in row:
                            coord = getattr(cell, "coordinate", "")
                            formula = formula_text.get((sheet_name, coord))
                            if formula:
                                formula_lines.append(f"{coord}: {formula}")
                            if cell.value is None:
                                continue
                            row_values.append(str(cell.value))
                            fmt = cell.number_format
                            # Excel applies a number format only to numeric cells; a
                            # text/bool value shows raw, so don't tag it as formatted.
                            numeric = isinstance(
                                cell.value, int | float
                            ) and not isinstance(cell.value, bool)
                            if fmt and numeric and _is_numeric_display_format(fmt):
                                # Emit the [fmt] tag; prepend the rendered value only
                                # when we can reproduce it.
                                rendered = _render_number_format(cell.value, fmt)
                                body = f"{rendered} " if rendered else ""
                                format_lines.append(
                                    f"{cell.coordinate}: {body}[fmt: {fmt}]"
                                )
                        if row_values:
                            sheet_text_lines.append("\t".join(row_values))

                    sheet_text = "\n".join(sheet_text_lines)
                    if format_lines:
                        sheet_text += "\n\n=== Number Formats ===\n" + "\n".join(
                            format_lines
                        )
                    if formula_lines:
                        sheet_text += "\n\n=== Formulas ===\n" + "\n".join(
                            formula_lines
                        )

                    sheet_text = f"=== Sheet: {sheet_name} ===\n{sheet_text}"

                    # Banner also goes on each sheet: graders may score a single
                    # sheet, which wouldn't include the parent text that carries it.
                    sub_artifacts.append(
                        SubArtifact(
                            index=sheet_idx,
                            type="sheet",
                            title=sheet_name,
                            content=unresolved_banner + sheet_text,
                            images=[],
                        )
                    )

                    if sub_artifact_index is None:
                        full_text_parts.append(sheet_text)
            finally:
                wb.close()

            # If specific sub-artifact requested, return only that (skip chart extraction)
            if sub_artifact_index is not None:
                if sub_artifacts:
                    return ExtractedContent(
                        # content already carries unresolved_banner (see above)
                        text=sub_artifacts[0].content,
                        images=[],
                        extraction_method=self.name,
                        metadata={"sheet_index": sub_artifact_index},
                        sub_artifacts=[],
                    )
                elif skipped_hidden_sheet:
                    raise ValueError(
                        f"Sheet index {sub_artifact_index} ('{skipped_hidden_sheet}') is hidden"
                    )
                else:
                    raise ValueError(f"Sheet index {sub_artifact_index} not found")

            # Extract charts if present (only for full file extraction)
            chart_images: list[ImageMetadata] = []
            if has_charts_in_xlsx(file_path):
                logger.info(f"Charts detected in {file_path.name}")

                soffice_path = find_libreoffice()
                if soffice_path:
                    pdf_path = await convert_xlsx_to_pdf(file_path, soffice_path)
                    # Start try immediately to ensure cleanup on CancelledError
                    try:
                        if pdf_path:
                            chart_images = pdf_to_base64_images(pdf_path)
                            if chart_images:
                                logger.info(
                                    f"Extracted {len(chart_images)} chart image(s) from PDF"
                                )

                                # Add chart placeholders to text
                                chart_text = "\n\n=== Charts ===\n"
                                for img in chart_images:
                                    chart_text += f"{img.placeholder} - {img.caption}\n"
                                full_text_parts.append(chart_text)
                    finally:
                        if pdf_path:
                            if pdf_path.exists():
                                pdf_path.unlink()
                            if pdf_path.parent.exists():
                                shutil.rmtree(pdf_path.parent, ignore_errors=True)
                else:
                    logger.warning(
                        f"LibreOffice not found - cannot extract chart images from {file_path.name}. "
                        "Install LibreOffice for chart extraction support."
                    )

            # Return all sheets with chart images
            full_text = unresolved_banner + "\n\n".join(full_text_parts)
            return ExtractedContent(
                text=full_text,
                images=chart_images,
                extraction_method=self.name,
                metadata={
                    "sheet_count": len(sub_artifacts),
                    "chart_count": len(chart_images),
                },
                sub_artifacts=sub_artifacts,
            )

        except Exception as e:
            logger.warning(f"Failed to extract XLSX with openpyxl: {e}")
            raise

    @retry(stop=stop_after_attempt(3), wait=wait_fixed(1), reraise=True)
    async def _extract_xls(
        self, file_path: Path, sub_artifact_index: int | None = None
    ) -> ExtractedContent:
        """Extract content from XLS file using xlrd"""
        if xlrd is None:
            raise ImportError("xlrd is required for XLS extraction")

        try:
            wb = xlrd.open_workbook(str(file_path))
            sub_artifacts = []
            full_text_parts = []

            for sheet_idx in range(wb.nsheets):
                if sub_artifact_index is not None and sheet_idx != sub_artifact_index:
                    continue

                sheet = wb.sheet_by_index(sheet_idx)
                sheet_name = sheet.name

                # Skip hidden sheets (visibility: 0=visible, 1=hidden, 2=very hidden)
                if sheet.visibility != 0:
                    continue

                sheet_text_lines = []
                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.ctype == xlrd.XL_CELL_EMPTY:
                            continue
                        elif cell.ctype == xlrd.XL_CELL_NUMBER:
                            try:
                                value = cell.value
                                # Check for special float values (NaN, inf, -inf)
                                if isinstance(value, float) and (
                                    math.isnan(value) or math.isinf(value)
                                ):
                                    row_values.append(str(value))
                                elif value == int(value):
                                    row_values.append(str(int(value)))
                                else:
                                    row_values.append(str(value))
                            except (ValueError, OverflowError, TypeError):
                                row_values.append(str(cell.value))
                        elif cell.ctype == xlrd.XL_CELL_DATE:
                            try:
                                dt = xlrd.xldate_as_tuple(
                                    float(cell.value), wb.datemode
                                )
                                row_values.append(f"{dt[0]}-{dt[1]:02d}-{dt[2]:02d}")
                            except Exception:
                                row_values.append(str(cell.value))
                        elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                            row_values.append("TRUE" if cell.value else "FALSE")
                        else:
                            value = str(cell.value).strip()
                            if value:
                                row_values.append(value)

                    if row_values:
                        sheet_text_lines.append("\t".join(row_values))

                sheet_text = "\n".join(sheet_text_lines)
                sheet_text = f"=== Sheet: {sheet_name} ===\n{sheet_text}"

                sub_artifacts.append(
                    SubArtifact(
                        index=sheet_idx,
                        type="sheet",
                        title=sheet_name,
                        content=sheet_text,
                        images=[],
                    )
                )

                if sub_artifact_index is None:
                    full_text_parts.append(sheet_text)

            logger.debug(
                f"[LOCAL] Extracted {len(sub_artifacts)} sub-artifacts from {file_path}"
            )

            if sub_artifact_index is not None:
                if sub_artifacts:
                    return ExtractedContent(
                        text=sub_artifacts[0].content,
                        images=[],
                        extraction_method=self.name,
                        metadata={"sheet_index": sub_artifact_index},
                        sub_artifacts=[],
                    )
                else:
                    raise ValueError(f"Sheet index {sub_artifact_index} not found")

            return ExtractedContent(
                text="\n\n".join(full_text_parts),
                images=[],
                extraction_method=self.name,
                metadata={"sheet_count": len(sub_artifacts)},
                sub_artifacts=sub_artifacts,
            )

        except Exception as e:
            logger.warning(f"Failed to extract XLS with xlrd: {e}")
            raise

    def _extract_text_from_shape(self, shape: Any) -> list[str]:
        """
        Recursively extract text from a PowerPoint shape.

        Handles:
        - Simple shapes with .text attribute
        - Tables (extracts all cells)
        - Grouped shapes (recursively extracts from children)
        - Text frames with paragraphs
        """
        text_parts = []

        # Handle grouped shapes recursively
        if hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                text_parts.extend(self._extract_text_from_shape(child_shape))
            return text_parts

        # Handle tables - try to extract, but fall through if not a table
        # Note: hasattr(shape, "table") returns True for all GraphicFrame shapes
        # (charts, diagrams, etc.), but .table raises ValueError for non-tables
        try:
            table = shape.table
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    text_parts.append("\t".join(row_texts))
            return text_parts  # Only return if table extraction succeeded
        except (ValueError, AttributeError):
            pass  # Not a table shape, continue with other extraction methods

        # Handle text frames (more thorough than just .text)
        text_frame_succeeded = False
        if hasattr(shape, "text_frame"):
            try:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    para_text = ""
                    for run in paragraph.runs:
                        if run.text:
                            para_text += run.text
                    if para_text.strip():
                        text_parts.append(para_text.strip())
                        text_frame_succeeded = True
            except Exception:
                pass

        # Fallback to simple .text attribute (runs if text_frame failed or found nothing)
        if not text_frame_succeeded and hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                text_parts.append(text)

        return text_parts

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_fixed(1),
        reraise=True,
    )
    async def _extract_pptx(
        self, file_path: Path, sub_artifact_index: int | None = None
    ) -> ExtractedContent:
        """Extract content from PPTX file using python-pptx"""
        if Presentation is None:
            raise ImportError("python-pptx is required for PPTX extraction")

        try:
            prs = Presentation(str(file_path))

            sub_artifacts = []
            full_text_parts = []

            for slide_idx, slide in enumerate(prs.slides):
                # If specific sub-artifact requested, skip others
                if sub_artifact_index is not None and slide_idx != sub_artifact_index:
                    continue

                # Extract text from all shapes in the slide (including tables, groups, etc.)
                slide_text_parts = []
                slide_title = None

                for shape in slide.shapes:
                    # Try to detect title placeholder first
                    if slide_title is None:
                        try:
                            if (
                                hasattr(shape, "placeholder_format")
                                and shape.placeholder_format.type == 1
                            ):
                                shape_text = getattr(shape, "text", None)
                                if shape_text:
                                    slide_title = shape_text.strip()
                        except Exception:
                            pass

                    # Extract all text from this shape (recursively handles tables, groups, etc.)
                    shape_texts = self._extract_text_from_shape(shape)
                    slide_text_parts.extend(shape_texts)

                slide_text = "\n".join(slide_text_parts)

                # Use first line as title if no title detected
                if slide_title is None and slide_text_parts:
                    slide_title = slide_text_parts[0][:100]  # First 100 chars

                # Create sub-artifact for this slide
                sub_artifacts.append(
                    SubArtifact(
                        index=slide_idx,
                        type="slide",
                        title=slide_title or f"Slide {slide_idx + 1}",
                        content=slide_text,
                        images=[],
                    )
                )

                # Add to full text if not requesting specific sub-artifact
                if sub_artifact_index is None:
                    full_text_parts.append(
                        f"=== Slide {slide_idx + 1}: {slide_title or 'Untitled'} ===\n{slide_text}"
                    )

            # If specific sub-artifact requested, return only that
            if sub_artifact_index is not None:
                if sub_artifacts:
                    return ExtractedContent(
                        text=sub_artifacts[0].content,
                        images=[],
                        extraction_method=self.name,
                        metadata={"slide_index": sub_artifact_index},
                        sub_artifacts=[],  # Empty list when extracting single sub-artifact
                    )
                else:
                    raise ValueError(f"Slide index {sub_artifact_index} not found")

            # Return all slides
            full_text = "\n\n".join(full_text_parts)
            return ExtractedContent(
                text=full_text,
                images=[],
                extraction_method=self.name,
                metadata={"slide_count": len(sub_artifacts)},
                sub_artifacts=sub_artifacts,
            )

        except Exception as e:
            logger.warning(f"Failed to extract PPTX with python-pptx: {e}")
            raise

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_fixed(1),
        reraise=True,
    )
    async def _extract_docx(
        self, file_path: Path, sub_artifact_index: int | None = None
    ) -> ExtractedContent:
        """
        Extract content from DOCX file using python-docx.

        Note: python-docx doesn't have page concept, so we extract sections or the full document.
        For page-level extraction, Reducto is used when changes are detected.
        """
        if Document is None:
            raise ImportError("python-docx is required for DOCX extraction")

        try:
            doc = Document(str(file_path))

            comment_anchors = _build_comment_anchor_map(doc)

            # Extract all paragraphs
            all_text_parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    all_text_parts.append(text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        # Join cells with tabs to preserve table structure
                        all_text_parts.append("\t".join(row_text))

            # Extract document comments
            if len(doc.comments) > 0:
                all_text_parts.append("\n=== DOCUMENT COMMENTS ===")
                for comment in doc.comments:
                    header_parts = [f"Author: {comment.author}"]
                    if comment.timestamp:
                        header_parts.append(
                            f"Date: {comment.timestamp.strftime('%Y-%m-%d %H:%M')}"
                        )
                    cid = comment._element.get(qn("w:id"))  # noqa: SLF001
                    anchor = comment_anchors.get(cid) if cid is not None else None
                    if anchor:
                        anchored_text = anchor["anchored_text"]
                        location_line = (
                            f"Anchored to paragraph {anchor['paragraph_index']}"
                            + (
                                f': "{anchored_text}"'
                                if anchored_text
                                else " (zero-length range)"
                            )
                        )
                    else:
                        location_line = "Anchor not found in document body"
                    all_text_parts.append(
                        f"[Comment] ({', '.join(header_parts)})\n"
                        f"  {location_line}\n"
                        f"  Comment text: {comment.text}"
                    )

            full_text = "\n".join(all_text_parts)

            # For local extraction, we treat the whole document as one unit for change detection
            # We don't create sub-artifacts here because python-docx doesn't have reliable page info
            # If changes are detected, Reducto will handle page-level extraction

            # Return as single artifact (no sub-artifacts for simple change detection)
            return ExtractedContent(
                text=full_text,
                images=[],
                extraction_method=self.name,
                metadata={
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                    "comment_count": len(doc.comments),
                },
                sub_artifacts=[],  # No sub-artifacts - will use Reducto if changes detected
            )

        except Exception as e:
            logger.warning(f"Failed to extract DOCX with python-docx: {e}")
            raise

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_fixed(1),
        reraise=True,
    )
    async def _extract_pdf(
        self, file_path: Path, sub_artifact_index: int | None = None
    ) -> ExtractedContent:
        """Extract text from a PDF with pypdf (offline last-resort tier).

        Text only — no OCR or layout — but keeps PDFs readable when
        Reducto/Mercor (preferred, tried first) are unconfigured or fail.
        """
        if PdfReader is None:
            raise ImportError("pypdf is required for PDF extraction")

        try:
            reader = PdfReader(str(file_path))
            # Empty/owner-password PDFs open after decrypt(""); a truly locked one
            # raises later at reader.pages, erroring the trial.
            if reader.is_encrypted:
                reader.decrypt("")

            sub_artifacts = []
            full_text_parts = []

            for page_idx, page in enumerate(reader.pages):
                # If specific sub-artifact requested, skip others
                if sub_artifact_index is not None and page_idx != sub_artifact_index:
                    continue

                try:
                    page_text = (page.extract_text() or "").strip()
                except Exception as e:
                    logger.warning(
                        f"Failed to extract text from page {page_idx + 1} of "
                        f"{file_path.name}: {e}"
                    )
                    # Marker (not "") so a failed page stays visible and can't
                    # masquerade as blank in an otherwise-complete document.
                    page_text = (
                        f"[page {page_idx + 1} text extraction failed: "
                        f"{type(e).__name__}]"
                    )

                sub_artifacts.append(
                    SubArtifact(
                        index=page_idx,
                        type="page",
                        title=f"Page {page_idx + 1}",
                        content=page_text,
                        images=[],
                    )
                )

                if sub_artifact_index is None:
                    full_text_parts.append(f"=== Page {page_idx + 1} ===\n{page_text}")

            # Whole document with no text is likely a scan (no OCR offline); flag it.
            # Decide on the full doc only — a single requested page can't tell.
            scanned_banner = ""
            if (
                sub_artifact_index is None
                and sub_artifacts
                and not any(sa.content.strip() for sa in sub_artifacts)
            ):
                logger.warning(
                    f"[LOCAL] {file_path.name}: {len(sub_artifacts)} page(s) but no "
                    "extractable text (likely scanned/image-only — no OCR)."
                )
                scanned_banner = (
                    "=== Extraction Note ===\n[no extractable text — the PDF is "
                    "likely scanned/image-only and no OCR is available offline]\n\n"
                )
                for sa in sub_artifacts:
                    sa.content = scanned_banner + sa.content

            # If specific sub-artifact requested, return only that
            if sub_artifact_index is not None:
                if sub_artifacts:
                    return ExtractedContent(
                        text=sub_artifacts[0].content,
                        images=[],
                        extraction_method=self.name,
                        metadata={"page_index": sub_artifact_index},
                        sub_artifacts=[],
                    )
                else:
                    raise ValueError(f"Page index {sub_artifact_index} not found")

            # Return all pages
            full_text = scanned_banner + "\n\n".join(full_text_parts)
            return ExtractedContent(
                text=full_text,
                images=[],
                extraction_method=self.name,
                metadata={"page_count": len(sub_artifacts)},
                sub_artifacts=sub_artifacts,
            )

        except Exception as e:
            logger.warning(f"Failed to extract PDF with pypdf: {e}")
            raise

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_fixed(1),
        reraise=True,
    )
    async def _extract_csv(self, file_path: Path) -> ExtractedContent:
        """
        Extract content from CSV file using built-in csv module.

        CSV files are treated as single artifacts (no sub-artifacts).
        This provides a fallback when Reducto fails (e.g., file too large).
        """
        try:
            # Try UTF-8 first, fallback to other encodings
            encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
            content_lines = None
            used_encoding = None

            for encoding in encodings:
                try:
                    with open(file_path, encoding=encoding, newline="") as f:
                        reader = csv.reader(f)
                        content_lines = []
                        for row in reader:
                            # Join cells with tabs to preserve structure
                            content_lines.append("\t".join(row))
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    content_lines = None
                    continue

            if content_lines is None:
                raise ValueError("Could not decode CSV with any supported encoding")

            full_text = "\n".join(content_lines)

            logger.debug(
                f"Extracted CSV with {len(content_lines)} rows using {used_encoding} encoding"
            )

            return ExtractedContent(
                text=full_text,
                images=[],
                extraction_method=self.name,
                metadata={
                    "row_count": len(content_lines),
                    "encoding": used_encoding,
                },
                sub_artifacts=[],  # CSV is treated as a single artifact
            )

        except Exception as e:
            logger.warning(f"Failed to extract CSV: {e}")
            raise
