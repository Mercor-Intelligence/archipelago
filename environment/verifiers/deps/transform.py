# flake8: noqa
#!/usr/bin/env python3
"""GDM Docker World artifact transformation script.

Standalone script packaged into the GDM export ZIP. Runs inside the
GDM container at shutdown time to transform reference artifacts into
grading-ready files.

Only used by OUTPUT_LLM_MULTI_REPRESENTATION verifiers. Other eval types
continue using simple file copies (mkdir + cp).

Supported transformations:
  *_to_text           - Extract text content from PDF/DOCX/PPTX
  *_to_images         - Render pages as a single concatenated PNG image (PDF/DOCX/PPTX only)
  *_to_pdf            - Convert Office docs to PDF via LibreOffice
  *_to_text_formulas  - Extract spreadsheet content preserving raw formulas
  spreadsheet_to_pdf  - Convert spreadsheet to PDF (with sheet names on each page)

Native (*_native) transformations are NOT handled here. They are treated
identically to source (just cp) by the GDM export pipeline.

Usage:
    python3 transform.py --input /app/files/report.docx \
                          --output /app/output/report.pdf \
                          --transformation docx_to_pdf
"""

from __future__ import annotations

import argparse
import io
import shutil
import subprocess
import sys
import tempfile
from collections.abc import Callable
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from PIL import Image as PILImage

SUPPORTED_TRANSFORMATIONS = {"docx_to_images", "docx_to_pdf", "docx_to_text", "pdf_to_images", "pdf_to_text", "pptx_to_images", "pptx_to_pdf", "pptx_to_text", "spreadsheet_to_pdf", "spreadsheet_to_text_formulas"}

LIBREOFFICE_TIMEOUT = 120
MAX_RENDER_PAGES = 50
RENDER_DPI = 150


def _find_libreoffice() -> str | None:
    return shutil.which("libreoffice") or shutil.which("soffice")


def _convert_to_pdf(
    input_path: Path,
    *,
    export_filter: str = "pdf",
    profile_xml: str | None = None,
) -> Path | None:
    soffice = _find_libreoffice()
    if not soffice:
        print("ERROR: LibreOffice not found", file=sys.stderr)
        return None

    out_dir = tempfile.mkdtemp(prefix="transform_pdf_")
    user_dir = tempfile.mkdtemp(prefix="libreoffice_profile_")

    if profile_xml:
        profile_user = Path(user_dir) / "user"
        profile_user.mkdir(parents=True, exist_ok=True)
        (profile_user / "registrymodifications.xcu").write_text(
            profile_xml, encoding="utf-8"
        )

    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--nologo",
                "--nolockcheck",
                f"-env:UserInstallation=file://{user_dir}",
                "--convert-to",
                export_filter,
                "--outdir",
                out_dir,
                str(input_path),
            ],
            capture_output=True,
            timeout=LIBREOFFICE_TIMEOUT,
        )

        if result.returncode != 0:
            print(
                f"ERROR: LibreOffice conversion failed: {result.stderr.decode()}",
                file=sys.stderr,
            )
            shutil.rmtree(out_dir, ignore_errors=True)
            return None

        pdf_path = Path(out_dir) / f"{input_path.stem}.pdf"
        if pdf_path.exists():
            return pdf_path

        print(f"ERROR: PDF not found at {pdf_path}", file=sys.stderr)
        shutil.rmtree(out_dir, ignore_errors=True)
        return None

    except subprocess.TimeoutExpired:
        print("ERROR: LibreOffice conversion timed out", file=sys.stderr)
        shutil.rmtree(out_dir, ignore_errors=True)
        return None
    finally:
        shutil.rmtree(user_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------


def _extract_text_pdf(input_path: Path) -> str:
    import fitz  # pyright: ignore[reportMissingImports]

    doc = fitz.open(str(input_path))
    pages = []
    try:
        for page in doc:
            pages.append(page.get_text())
    finally:
        doc.close()
    return "\n".join(pages)


def _extract_text_docx(input_path: Path) -> str:
    from docx import Document

    doc = Document(str(input_path))
    paragraphs = [p.text for p in doc.paragraphs]

    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append("\t".join(cell.text for cell in row.cells))
        paragraphs.append("\n".join(rows))

    return "\n".join(paragraphs)


def _extract_text_pptx(input_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(input_path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:  # pyright: ignore[reportAttributeAccessIssue]
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if hasattr(shape, "table") and shape.has_table:
                for row in shape.table.rows:  # pyright: ignore[reportAttributeAccessIssue]
                    row_text = "\t".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        slide_texts.append(row_text)
        if slide_texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _col_letter(col_idx: int) -> str:
    """Convert 1-based column index to Excel letter (1=A, 27=AA, etc.)."""
    result = ""
    while col_idx > 0:
        col_idx, rem = divmod(col_idx - 1, 26)
        result = chr(65 + rem) + result
    return result


def _xml_escape(val: str) -> str:
    return val.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


def _spreadsheet_to_xml(input_path: Path, *, data_only: bool) -> str:
    import openpyxl

    read_only = data_only
    wb = openpyxl.load_workbook(str(input_path), data_only=data_only, read_only=read_only)
    parts = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_xml: list[str] = []
            for row_idx, row in enumerate(ws.iter_rows(min_row=1), start=1):
                cells_xml: list[str] = []
                has_value = False
                for cell in row:
                    col = _col_letter(cell.column) if hasattr(cell, "column") else _col_letter(cells_xml.__len__() + 1)
                    ref = f"{col}{row_idx}"
                    val = cell.value
                    if val is None:
                        cells_xml.append(f'    <cell ref="{ref}" />')
                    else:
                        has_value = True
                        escaped = _xml_escape(str(val))
                        if isinstance(val, str) and not str(val).startswith("="):
                            cells_xml.append(f'    <cell ref="{ref}" type="string">{escaped}</cell>')
                        else:
                            cells_xml.append(f'    <cell ref="{ref}">{escaped}</cell>')
                if has_value:
                    rows_xml.append(
                        f'  <row number="{row_idx}">\n' + "\n".join(cells_xml) + "\n  </row>"
                    )
            if rows_xml:
                escaped_name = _xml_escape(sheet_name)
                parts.append(
                    f'<template sheet="{escaped_name}">\n' + "\n".join(rows_xml) + "\n</template>"
                )
    finally:
        wb.close()
    return "\n\n".join(parts)


def _extract_text_spreadsheet_formulas(input_path: Path) -> str:
    return _spreadsheet_to_xml(input_path, data_only=False)


_TEXT_EXTRACTORS: dict[str, Callable[[Path], str]] = {
    "pdf_to_text": _extract_text_pdf,
    "docx_to_text": _extract_text_docx,
    "pptx_to_text": _extract_text_pptx,
    "spreadsheet_to_text_formulas": _extract_text_spreadsheet_formulas,
}


# ---------------------------------------------------------------------------
# PDF conversion
# ---------------------------------------------------------------------------

_CALC_BACKGROUND_PROFILE_XML = """\
<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry" \
xmlns:xs="http://www.w3.org/2001/XMLSchema" \
xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
<item oor:path="/org.openoffice.Office.Calc/Print/Page">\
<prop oor:name="Background" oor:op="fuse"><value>true</value></prop></item>
<item oor:path="/org.openoffice.Office.Calc/Print/Page">\
<prop oor:name="BlackWhite" oor:op="fuse"><value>false</value></prop></item>
</oor:items>"""


def _convert_docx_to_pdf(input_path: Path, output_path: Path) -> None:
    pdf = _convert_to_pdf(input_path, export_filter="pdf:writer_pdf_Export")
    if not pdf:
        print(f"ERROR: Failed to convert {input_path} to PDF", file=sys.stderr)
        sys.exit(1)
    try:
        shutil.move(str(pdf), str(output_path))
    finally:
        if pdf.parent.exists():
            shutil.rmtree(pdf.parent, ignore_errors=True)


def _convert_pptx_to_pdf(input_path: Path, output_path: Path) -> None:
    pdf = _convert_to_pdf(input_path, export_filter="pdf:impress_pdf_Export")
    if not pdf:
        print(f"ERROR: Failed to convert {input_path} to PDF", file=sys.stderr)
        sys.exit(1)
    try:
        shutil.move(str(pdf), str(output_path))
    finally:
        if pdf.parent.exists():
            shutil.rmtree(pdf.parent, ignore_errors=True)


def _convert_spreadsheet_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert spreadsheet to PDF with "Sheet: {name}" on every page.

    Converts each sheet separately so we know exactly which sheet each PDF
    page came from, including wide/long sheets that span multiple pages.
    """
    import fitz  # pyright: ignore[reportMissingImports]
    import openpyxl

    wb = openpyxl.load_workbook(str(input_path))
    sheet_names = list(wb.sheetnames)
    wb.close()

    if not sheet_names:
        print(f"ERROR: No sheets in {input_path}", file=sys.stderr)
        sys.exit(1)

    out_dir = tempfile.mkdtemp(prefix="transform_spreadsheet_pdf_")
    merged_path = Path(out_dir) / "merged.pdf"
    try:
        with fitz.open() as merged_doc:
            for sheet_name in sheet_names:
                # Create temp workbook with only this sheet
                wb_one = openpyxl.load_workbook(str(input_path))
                for name in list(wb_one.sheetnames):
                    if name != sheet_name:
                        del wb_one[name]
                wb_one.active = wb_one[sheet_name]
                sheet_xlsx = Path(out_dir) / f"sheet_{_safe_filename(sheet_name)}.xlsx"
                wb_one.save(str(sheet_xlsx))
                wb_one.close()

                pdf = _convert_to_pdf(
                    sheet_xlsx,
                    export_filter="pdf:calc_pdf_Export",
                    profile_xml=_CALC_BACKGROUND_PROFILE_XML,
                )
                if not pdf:
                    print(
                        f"ERROR: Failed to convert sheet '{sheet_name}' to PDF",
                        file=sys.stderr,
                    )
                    sys.exit(1)

                try:
                    with fitz.open(str(pdf)) as sheet_doc:
                        for page in sheet_doc:
                            label = f"Sheet: {sheet_name}"
                            page.insert_text(
                                (36, 18), label, fontsize=9, fontname="Helvetica"
                            )
                        merged_doc.insert_pdf(sheet_doc)
                finally:
                    if pdf.exists():
                        pdf.unlink(missing_ok=True)
                    if pdf.parent.exists():
                        shutil.rmtree(pdf.parent, ignore_errors=True)
                sheet_xlsx.unlink(missing_ok=True)

            merged_doc.save(str(merged_path))
        shutil.move(str(merged_path), str(output_path))
    finally:
        if Path(out_dir).exists():
            shutil.rmtree(out_dir, ignore_errors=True)


def _safe_filename(s: str) -> str:
    """Replace characters unsafe for temp filenames."""
    return "".join(c if c.isalnum() or c in "._-" else "_" for c in s)[:50]


_PDF_CONVERTERS: dict[str, Callable[[Path, Path], None]] = {
    "docx_to_pdf": _convert_docx_to_pdf,
    "pptx_to_pdf": _convert_pptx_to_pdf,
    "spreadsheet_to_pdf": _convert_spreadsheet_to_pdf,
}


# ---------------------------------------------------------------------------
# Image rendering
# ---------------------------------------------------------------------------


def _render_pdf_pages(pdf_path: Path) -> list[PILImage.Image]:
    import fitz  # pyright: ignore[reportMissingImports]
    from PIL import Image

    doc = fitz.open(str(pdf_path))
    images: list[PILImage.Image] = []
    try:
        for page in doc[:MAX_RENDER_PAGES]:  # pyright: ignore[reportArgumentType]
            zoom = RENDER_DPI / 72.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            images.append(img)
    finally:
        doc.close()
    return images


def _concatenate_images(images: list[PILImage.Image]) -> bytes:
    from PIL import Image

    if not images:
        return b""
    if len(images) == 1:
        buf = io.BytesIO()
        images[0].save(buf, format="PNG")
        return buf.getvalue()

    total_width = max(img.width for img in images)
    total_height = sum(img.height for img in images)
    combined = Image.new("RGB", (total_width, total_height), (255, 255, 255))
    y_offset = 0
    for img in images:
        combined.paste(img, (0, y_offset))
        y_offset += img.height

    buf = io.BytesIO()
    combined.save(buf, format="PNG")
    return buf.getvalue()


def _render_to_images(input_path: Path, transformation: str) -> bytes:
    ext = input_path.suffix.lower()
    needs_libreoffice = ext != ".pdf"

    if needs_libreoffice:
        pdf_path = _convert_to_pdf(input_path)
        if not pdf_path:
            print(
                f"ERROR: Failed to convert {input_path} to PDF",
                file=sys.stderr,
            )
            sys.exit(1)
        try:
            pages = _render_pdf_pages(pdf_path)
        finally:
            if pdf_path.exists():
                pdf_path.unlink(missing_ok=True)
            if pdf_path.parent.exists():
                shutil.rmtree(pdf_path.parent, ignore_errors=True)
    else:
        pages = _render_pdf_pages(input_path)

    if not pages:
        print(f"ERROR: No pages rendered from {input_path}", file=sys.stderr)
        sys.exit(1)

    return _concatenate_images(pages)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Transform a reference artifact for GDM grading."
    )
    parser.add_argument("--input", required=True, help="Path to the input file")
    parser.add_argument("--output", required=True, help="Path for the output file")
    parser.add_argument(
        "--transformation",
        required=True,
        choices=sorted(SUPPORTED_TRANSFORMATIONS),
        help="Transformation to apply",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output)
    transformation = args.transformation

    if not input_path.exists():
        print(f"ERROR: Input file not found: {input_path}", file=sys.stderr)
        return 1

    output_path.parent.mkdir(parents=True, exist_ok=True)

    if transformation in _TEXT_EXTRACTORS:
        extractor = _TEXT_EXTRACTORS[transformation]
        text = extractor(input_path)
        output_path.write_text(text, encoding="utf-8")
        print(
            f"OK: {transformation} | {input_path} -> {output_path} "
            f"({len(text):,} chars)",
            file=sys.stderr,
        )
    elif transformation in _PDF_CONVERTERS:
        converter = _PDF_CONVERTERS[transformation]
        converter(input_path, output_path)
        size = output_path.stat().st_size
        print(
            f"OK: {transformation} | {input_path} -> {output_path} "
            f"({size:,} bytes)",
            file=sys.stderr,
        )
    else:
        image_bytes = _render_to_images(input_path, transformation)
        output_path.write_bytes(image_bytes)
        print(
            f"OK: {transformation} | {input_path} -> {output_path} "
            f"({len(image_bytes):,} bytes)",
            file=sys.stderr,
        )

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
