import os
from io import BytesIO

from models.response import ReadDeckResponse, SlideOverviewData
from models.tool_inputs import ReadCompleteDeckInput
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from utils.decorators import make_async_background
from utils.path_utils import resolve_under_root

SLIDES_ROOT = os.getenv("APP_SLIDES_ROOT") or os.getenv("APP_FS_ROOT", "/filesystem")


def _resolve_under_root(path: str) -> str:
    """Map path to the slides root, refusing anything that escapes it.

    Delegates to the shared resolver so containment is enforced in ONE place:
    the previous per-module copy normpath()'d after joining and never checked
    the result, so ``/../app/tools/...`` resolved outside the root. Kept as a
    module-level name so call sites (and tests) are unchanged.

    Raises PathTraversalError when the path escapes the sandbox.
    """
    return resolve_under_root(path, root=SLIDES_ROOT)


@make_async_background
def read_completedeck(request: ReadCompleteDeckInput) -> ReadDeckResponse:
    """Read overview information about all slides in a PowerPoint presentation.

    Provides a high-level summary of the entire deck with slide titles and content previews.
    Ideal for understanding deck structure before performing detailed operations.

    Notes:
        - Fast overview without detailed components
        - Use slide_index to target slides with edit_slides/read_individualslide
        - For details: read_individualslide. For markdown: read_slides
        - Content preview: up to 500 chars, truncated without indicator
        - Title defaults to 'Slide {index}' if no text found
        - Speaker notes excluded
    """

    def error(msg: str) -> ReadDeckResponse:
        return ReadDeckResponse(success=False, error=msg)

    target_path = _resolve_under_root(request.file_path)

    try:
        if not os.path.exists(target_path):
            return error(f"File not found: {request.file_path}")
        if not os.path.isfile(target_path):
            return error(f"Not a file: {request.file_path}")

        with open(target_path, "rb") as f:
            file_bytes = f.read()
    except Exception as exc:
        return error(f"Failed to read presentation: {repr(exc)}")

    try:
        presentation = Presentation(BytesIO(file_bytes))

        slides_data = []

        for slide_index, slide in enumerate(presentation.slides):
            # Extract title (usually from the first placeholder)
            title = ""
            content_parts = []

            for shape in slide.shapes:
                # Only process shapes that have text
                if not isinstance(shape, Shape) or not shape.has_text_frame:
                    continue

                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # Try to detect if this is the title
                    if (
                        not title
                        and hasattr(shape, "is_placeholder")
                        and shape.is_placeholder
                    ):
                        try:
                            if shape.placeholder_format.type in (0, 2):
                                title = text
                                continue
                        except (AttributeError, ValueError):
                            pass

                    if not title and not content_parts:
                        title = text
                    else:
                        content_parts.append(text)

            if not title:
                title = f"Slide {slide_index}"

            content = "\n".join(content_parts) if content_parts else "(No content)"

            slides_data.append(
                SlideOverviewData(
                    slide_index=slide_index,
                    title=title,
                    content=content,
                )
            )

        return ReadDeckResponse(
            success=True,
            total_slides=len(presentation.slides),
            slides=slides_data,
        )

    except Exception as exc:
        return error(f"Failed to parse presentation: {repr(exc)}")
